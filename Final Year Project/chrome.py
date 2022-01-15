# import subprocess, sys, time

# p = subprocess.Popen([sys.executable, 'AIVirtualMouseProject.py'], 
#                                     stdout=subprocess.PIPE, 
#                                     stderr=subprocess.STDOUT)
# print(1)
# time.sleep(15)
# p.kill()
from pptx import Presentation

filename = 'C:\\Users\\Anshir\\Desktop\\Final Year Project\\New PPT Presentation.pptx'

prs = Presentation(filename)
print(prs.slides)
for slide in prs.slides:
    #title = slide.shapes.title.text
    slide.get(0)